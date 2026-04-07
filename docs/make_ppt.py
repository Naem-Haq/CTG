from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path('/home/naem-haq/Software-Engineering/CTG')
OUT = ROOT / 'docs' / 'ctg_benchmark_presentation.pptx'
IMAGES = {
    'rf_cm': ROOT / 'outputs/models/rf_3class_confusion_matrix.png',
    'cnn_cm': ROOT / 'outputs/models/cnn_3class_confusion_matrix.png',
    'rf_fi': ROOT / 'outputs/models/rf_3class_feature_importance.png',
    'alerts': ROOT / 'outputs/models/blockF_alert_timeline.png',
}

PALETTE = {
    'bg': RGBColor(11, 16, 32),
    'text': RGBColor(236, 241, 255),
    'muted': RGBColor(185, 196, 226),
    'accent': RGBColor(62, 183, 219),
    'accent_2': RGBColor(51, 109, 217),
}


def _no_line(shape):
    shape.line.fill.background()


def _send_to_back(shape):
    sp_tree = shape._element.getparent()
    el = shape._element
    sp_tree.remove(el)
    sp_tree.insert(2, el)


def _style_paragraph(paragraph, size, color, bold=False):
    paragraph.font.size = Pt(size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    paragraph.font.name = 'Calibri'
    for run in paragraph.runs:
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = 'Calibri'


def add_background_art(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = PALETTE['bg']

    glow_left = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.2), Inches(-1.0), Inches(4.0), Inches(4.0))
    glow_left.fill.solid()
    glow_left.fill.fore_color.rgb = PALETTE['accent_2']
    glow_left.fill.transparency = 72
    _no_line(glow_left)
    _send_to_back(glow_left)

    glow_right = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(4.8), Inches(4.0), Inches(4.0))
    glow_right.fill.solid()
    glow_right.fill.fore_color.rgb = PALETTE['accent']
    glow_right.fill.transparency = 80
    _no_line(glow_right)
    _send_to_back(glow_right)

    top_band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.12))
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = PALETTE['accent']
    _no_line(top_band)
    _send_to_back(top_band)


def style_text_frame(tf, base_size=20, color=None):
    color = color or PALETTE['text']
    for p in tf.paragraphs:
        p.font.size = Pt(base_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'


def set_title(slide, title):
    slide.shapes.title.text = title
    tf = slide.shapes.title.text_frame
    _style_paragraph(tf.paragraphs[0], 34, PALETTE['text'], bold=True)


def add_bullets(prs, title, bullets, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_background_art(slide)
    set_title(slide, title)
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()

    if subtitle:
        p = tf.paragraphs[0]
        p.text = subtitle
        _style_paragraph(p, 20, PALETTE['accent'], bold=True)

    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if (subtitle or i > 0) else tf.paragraphs[0]
        p.text = b
        p.level = 0
        _style_paragraph(p, 20, PALETTE['text'])

    # subtle emphasis line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(1.35), Inches(2.2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = PALETTE['accent']
    _no_line(line)

    return slide


def add_two_images(prs, title, left_path, right_path, caption_left, caption_right):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_background_art(slide)
    set_title(slide, title)

    left_x, right_x = Inches(0.5), Inches(5.1)
    top_img = Inches(1.3)
    max_w, max_h = Inches(4.4), Inches(4.7)

    def place(path, x):
        with Image.open(path) as im:
            w, h = im.size
        ratio = min(max_w / w, max_h / h)
        disp_w = w * ratio
        disp_h = h * ratio
        y = top_img + (max_h - disp_h) / 2
        slide.shapes.add_picture(str(path), x, y, width=disp_w, height=disp_h)
        frame = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x - Inches(0.08),
            y - Inches(0.08),
            disp_w + Inches(0.16),
            disp_h + Inches(0.16),
        )
        frame.fill.background()
        frame.line.color.rgb = PALETTE['accent']
        frame.line.width = Pt(1.8)

    place(left_path, left_x)
    place(right_path, right_x)

    tb1 = slide.shapes.add_textbox(left_x, Inches(6.2), max_w, Inches(0.5)).text_frame
    tb1.text = caption_left
    _style_paragraph(tb1.paragraphs[0], 14, PALETTE['muted'])
    tb1.paragraphs[0].alignment = PP_ALIGN.CENTER

    tb2 = slide.shapes.add_textbox(right_x, Inches(6.2), max_w, Inches(0.5)).text_frame
    tb2.text = caption_right
    _style_paragraph(tb2.paragraphs[0], 14, PALETTE['muted'])
    tb2.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


def set_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = text


def build_presentation():
    prs = Presentation()

    cover = prs.slides.add_slide(prs.slide_layouts[0])
    add_background_art(cover)
    cover.shapes.title.text = 'Benchmarking ML Classifiers for CTG\nFetal Heart Rate Pattern Detection'
    _style_paragraph(cover.shapes.title.text_frame.paragraphs[0], 40, PALETTE['text'], bold=True)
    cover.placeholders[1].text = (
        'Naem Haq\n'
        'FYP Dissertation Presentation\n'
        'Residency 4 context: INFANT Research Centre (UCC)'
    )
    for p in cover.placeholders[1].text_frame.paragraphs:
        _style_paragraph(p, 18, PALETTE['muted'])

    badge = cover.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(5.95), Inches(4.2), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = PALETTE['accent_2']
    badge.fill.transparency = 20
    _no_line(badge)
    badge_tf = badge.text_frame
    badge_tf.text = 'CTU-UHB • RF vs 1D-CNN • Reproducible Benchmark'
    badge_tf.paragraphs[0].font.size = Pt(12)
    badge_tf.paragraphs[0].font.color.rgb = PALETTE['text']
    badge_tf.paragraphs[0].font.name = 'Calibri'
    badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    set_notes(
        cover,
        'Hi, I’m Naem Haq. This final year project benchmarks machine-learning classifiers for fetal heart rate pattern detection in Cardiotocography, CTG for short. The work was shaped during my 4th Residency at INFANT Research Centre, where interpretation variability in labour monitoring highlighted the need for controlled, reproducible evaluation.',
    )

    s2 = add_bullets(
        prs,
        'Motivation and Context',
        [
            'CTG is widely used in labour monitoring, but interpretation can be subjective.',
            'Clinical variability motivates transparent decision-support benchmarking.',
            'This benchmarking focus was shaped during Residency 4 at INFANT Research Centre.',
        ],
    )
    set_notes(
        s2,
        'CTG is a core intrapartum monitoring tool, but interpretation can be subjective and variable across '
        'clinicians and settings. That variability motivates computational decision support. During my Residency '
        '4 placement at INFANT, this practical need became clear, which directly motivated a controlled '
        'benchmarking study rather than an isolated model-building exercise.',
    )

    s3 = add_bullets(
        prs,
        'Aim and Scope',
        [
            'Aim: benchmark supervised classifiers for fetal heart rate pattern detection.',
            'Models compared: Random Forest (feature-based) vs 1D-CNN (raw-signal).',
            'Scope: within-cohort CTU-UHB benchmark with weakly supervised labels.',
            'Claim boundary: fixed-split comparative evidence, not inferential repeated-run ranking.',
        ],
    )
    set_notes(
        s3,
        'The aim was to benchmark two model families under the same conditions: Random Forest and 1D-CNN. '
        'This is a within-cohort CTU-UHB benchmark with weakly supervised labels. I keep claims bounded to '
        'fixed-split comparative evidence, so conclusions are transparent and proportional to executed analysis.',
    )

    s4 = add_bullets(
        prs,
        'Benchmark Design and Fairness Controls',
        [
            'Shared preprocessing and leakage-aware grouped splitting by record ID.',
            'Common evaluation criteria: class-wise and aggregate metrics.',
            'Additional evidence: efficiency, confusion analysis, threshold behaviour, interpretability.',
            'Reproducibility controls: fixed seed policy, split manifests, JSON/CSV artifacts.',
        ],
    )
    set_notes(
        s4,
        'A key contribution is fairness of comparison. Both models share the same preprocessing, grouped '
        'splitting, and evaluation framework. Record-level grouping prevents leakage between train and test '
        'windows from the same labour record. Reproducibility is supported by fixed seeds, stored split '
        'manifests, and exported artifacts for auditability.',
    )

    s5 = add_bullets(
        prs,
        'Pipeline',
        [
            'Signal quality control and preprocessing',
            'Windowing and feature extraction',
            'Model training with grouped holdout protocol',
            'Evaluation, error analysis, thresholding, and post-processing review',
        ],
    )
    set_notes(
        s5,
        'The pipeline runs from raw CTG through quality control, preprocessing, segmentation and feature '
        'construction, then model training and evaluation. Beyond headline metrics, I include class-wise error '
        'analysis, threshold behavior, and post-processing effects so the benchmark reflects operational '
        'decision-support needs.',
    )

    s6 = add_bullets(
        prs,
        'Headline Results',
        [
            'Held-out windows: RF 95.55% accuracy, macro-F1 0.945.',
            'Held-out windows: CNN 65.31% accuracy, macro-F1 0.584.',
            'Record-level CNN drops to 48.19% accuracy, macro-F1 0.383.',
            'Result: RF is the stronger current baseline under this data regime.',
        ],
    )
    set_notes(
        s6,
        'The performance gap is substantial on held-out windows: RF reaches 95.55% accuracy and macro-F1 '
        '0.945, while CNN reaches 65.31% and 0.584. At record level, CNN drops further to 48.19% accuracy '
        'and macro-F1 0.383. Under this tested regime, Random Forest is the stronger baseline.',
    )

    if IMAGES['rf_cm'].exists() and IMAGES['cnn_cm'].exists():
        s7 = add_two_images(
            prs,
            'Confusion Matrix Comparison',
            IMAGES['rf_cm'],
            IMAGES['cnn_cm'],
            'Random Forest: tighter class separation',
            '1D-CNN: greater off-diagonal confusion',
        )
        set_notes(
            s7,
            'These confusion matrices show how the aggregate metric gap appears structurally. RF displays '
            'stronger diagonal concentration, indicating better class separation. The CNN matrix shows broader '
            'off-diagonal confusion, especially around boundary classes, which contributes to weaker macro-F1 '
            'and lower stability after aggregation.',
        )

    if IMAGES['rf_fi'].exists() and IMAGES['alerts'].exists():
        s8 = add_two_images(
            prs,
            'Interpretability and Alert Dynamics',
            IMAGES['rf_fi'],
            IMAGES['alerts'],
            'RF feature/domain importance',
            'Temporal alert timeline after smoothing',
        )
        set_notes(
            s8,
            'On the left, RF importance outputs provide interpretable evidence about which domains drive model '
            'decisions, supporting transparency. On the right, the alert timeline illustrates temporal behavior '
            'under post-processing policy, showing that usability and alert stability depend on threshold and '
            'smoothing choices, not only classifier architecture.',
        )

    s9 = add_bullets(
        prs,
        'Conclusion',
        [
            'Under tested conditions, Random Forest is the more suitable transparent baseline.',
            'This does not imply deep learning is always worse in CTG.',
            'The contribution is an auditable benchmark reference for next-stage validation.',
        ],
    )
    set_notes(
        s9,
        'The conclusion is conditional and scoped: under the tested conditions, Random Forest is currently '
        'the more suitable transparent benchmark baseline. This is not a universal claim against deep learning. '
        'The study contributes an auditable foundation that future work can extend.',
    )

    s10 = add_bullets(
        prs,
        'Limitations and Next Steps',
        [
            'Findings are bounded to CTU-UHB and weak-label regime.',
            'Repeated grouped resampling and stronger uncertainty quantification are future extensions.',
            'External validation across broader datasets is required for stronger transportability claims.',
            'Pathway: prospective, clinician-facing translation aligned with INFANT/AI4Life goals.',
        ],
    )
    set_notes(
        s10,
        'These findings are specific to the CTU-UHB cohort and weak-label setting used in this benchmark. Next steps are uncertainty-focused repeat evaluation and external validation on broader datasets to support prospective clinician-facing translation in INFANT/AI4Life contexts.'
    )

    prs.save(OUT)
    print(OUT)


if __name__ == '__main__':
    build_presentation()
